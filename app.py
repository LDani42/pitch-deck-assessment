# app.py
import streamlit as st
import os
import json
import time
from dotenv import load_dotenv
import pandas as pd
import anthropic
import tempfile
import PyPDF2
import pptx
import docx
from io import BytesIO
import matplotlib.pyplot as plt
import seaborn as sns

# Load environment variables
load_dotenv()

# Initialize the Anthropic client
client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

# Define prompts for each Business Model Canvas component
BUSINESS_MODEL_CANVAS_PROMPTS = {
    "overview": """You are an expert pitch deck consultant. Analyze the provided pitch deck and give an overall assessment of the entire business model. 
    Focus on coherence between different components, overall viability, and how compelling the business case is.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "key_partners": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Key Partners section.
    Evaluate how well the pitch identifies strategic partners, supplier relationships, and the mutual value creation with these partners.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "key_activities": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Key Activities section.
    Evaluate how well the pitch articulates the most important activities the business must perform to make its business model work.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "key_resources": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Key Resources section.
    Evaluate how well the pitch describes the most important assets required to make the business model work.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "value_proposition": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Value Proposition section.
    Evaluate how well the pitch articulates the bundle of products and services that create value for specific customer segments.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "customer_relationships": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Customer Relationships section.
    Evaluate how well the pitch describes the types of relationships the company establishes with specific customer segments.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "channels": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Channels section.
    Evaluate how well the pitch describes how the company communicates with and reaches its customer segments to deliver the value proposition.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "customer_segments": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Customer Segments section.
    Evaluate how well the pitch defines the different groups of people or organizations the business aims to reach and serve.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "cost_structure": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Cost Structure section.
    Evaluate how well the pitch describes all costs incurred to operate the business model.
    Provide a score between 0-100 and specific, actionable feedback.""",
    
    "revenue_streams": """You are an expert pitch deck consultant. Analyze the provided pitch deck and assess the Revenue Streams section.
    Evaluate how well the pitch describes how the company generates cash from each customer segment and the overall pricing strategy.
    Provide a score between 0-100 and specific, actionable feedback."""
}

# Response format
RESPONSE_FORMAT = """
Respond in JSON format with the following structure:
{
  "score": [numerical score between 0-100],
  "feedback": "[detailed feedback with specific suggestions for improvement]"
}
"""

# Function to extract text from various file formats
def extract_text_from_file(uploaded_file):
    """Extract text content from PDF, PPTX, or DOCX files"""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        temp_file.write(uploaded_file.getvalue())
        temp_file_path = temp_file.name
    
    extracted_text = ""
    
    try:
        if file_extension == '.pdf':
            with open(temp_file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    extracted_text += f"\n--- Page {page_num + 1} ---\n"
                    extracted_text += page.extract_text()
        
        elif file_extension == '.pptx':
            presentation = pptx.Presentation(temp_file_path)
            for slide_num, slide in enumerate(presentation.slides):
                extracted_text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + "\n"
        
        elif file_extension == '.docx':
            doc = docx.Document(temp_file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
                
        elif file_extension == '.ppt':
            extracted_text = "PPT format has limited extraction capabilities. Consider converting to PPTX for better results."
        
        else:
            extracted_text = f"Unsupported file format: {file_extension}"
    
    finally:
        # Clean up temp file
        os.unlink(temp_file_path)
    
    return extracted_text

# Function to analyze a single component using Claude
def analyze_component(component, content):
    """Analyze a single component of the Business Model Canvas using Claude API"""
    prompt = BUSINESS_MODEL_CANVAS_PROMPTS[component]
    
    try:
        message = client.messages.create(
            model="claude-3-opus-20240229",
            max_tokens=1000,
            messages=[
                {
                    "role": "user",
                    "content": f"{prompt}\n\nHere is the pitch deck content to analyze:\n\n{content}\n\n{RESPONSE_FORMAT}"
                }
            ]
        )
        
        response_text = message.content[0].text
        
        # Try to parse JSON from the response
        try:
            # First, check if response is wrapped in markdown code blocks
            json_match = response_text.strip()
            if json_match.startswith('```json') and json_match.endswith('```'):
                json_match = json_match[7:-3].strip()
            
            # If it looks like JSON, parse it
            if json_match.startswith('{') and json_match.endswith('}'):
                return json.loads(json_match)
            
            # Fallback to regex extraction if needed
            import re
            score_match = re.search(r'"score"\s*:\s*(\d+)', response_text)
            feedback_match = re.search(r'"feedback"\s*:\s*"([^"]+)"', response_text)
            
            if score_match and feedback_match:
                return {
                    "score": int(score_match.group(1)),
                    "feedback": feedback_match.group(1)
                }
                
            # Last resort fallback
            return {
                "score": 50,
                "feedback": "Analysis could not be processed properly. Here's the raw response: " + response_text[:200] + "..."
            }
            
        except json.JSONDecodeError:
            # If JSON parsing fails, extract key information with regex
            import re
            score_match = re.search(r'"score"\s*:\s*(\d+)', response_text)
            feedback_match = re.search(r'"feedback"\s*:\s*"([^"]+)"', response_text)
            
            return {
                "score": int(score_match.group(1)) if score_match else 50,
                "feedback": feedback_match.group(1) if feedback_match else "Failed to parse Claude's response."
            }
            
    except Exception as e:
        st.error(f"Error analyzing {component}: {str(e)}")
        return {
            "score": 0,
            "feedback": f"Failed to analyze {component}. Please try again."
        }

# Main app
def main():
    st.set_page_config(
        page_title="Pitch Deck Assessment",
        page_icon="📊",
        layout="wide"
    )
    
    # App title and description
    st.title("Pitch Deck Assessment")
    st.markdown("Upload your pitch deck to get AI-powered feedback on your Business Model Canvas components.")
    
    # File uploader
    uploaded_file = st.file_uploader("Choose a file", type=["pdf", "pptx", "docx", "ppt"])
    
    if uploaded_file is not None:
        st.info(f"File uploaded: {uploaded_file.name}")
        
        # Add analyze button
        if st.button("Analyze Pitch Deck"):
            # Extract text from the file
            with st.spinner("Extracting text from the pitch deck..."):
                extracted_text = extract_text_from_file(uploaded_file)
                st.session_state.extracted_text = extracted_text
                
                # Display a snippet of extracted text for verification
                with st.expander("Preview of extracted text"):
                    st.text(extracted_text[:1000] + "..." if len(extracted_text) > 1000 else extracted_text)
            
            # Define components to analyze
            components = list(BUSINESS_MODEL_CANVAS_PROMPTS.keys())
            
            # Initialize results dictionary
            results = {}
            
            # Create a progress bar
            progress_bar = st.progress(0)
            progress_text = st.empty()
            
            # Analyze each component
            for i, component in enumerate(components):
                progress_percent = i / len(components)
                progress_bar.progress(progress_percent)
                progress_text.text(f"Analyzing {component.replace('_', ' ').title()}... ({i+1}/{len(components)})")
                
                # Analyze the component
                results[component] = analyze_component(component, extracted_text)
                
                # Small delay to update the UI
                time.sleep(0.1)
            
            # Complete the progress bar
            progress_bar.progress(1.0)
            progress_text.text("Analysis complete!")
            
            # Store results in session state
            st.session_state.results = results
            
            # Force a rerun to display results
            st.experimental_rerun()
    
    # Display results if they exist in session state
    if hasattr(st.session_state, 'results') and st.session_state.results:
        results = st.session_state.results
        
        # Create tabs for different views
        tab1, tab2 = st.tabs(["Overview Dashboard", "Detailed Analysis"])
        
        with tab1:
            st.header("Business Model Canvas Assessment")
            
            # Calculate average score
            avg_score = sum(component["score"] for component in results.values()) / len(results)
            
            # Display average score in a metric
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Average Score", f"{avg_score:.1f}/100")
            with col2:
                # Count strengths (scores >= 80)
                strengths = sum(1 for component in results.values() if component["score"] >= 80)
                st.metric("Key Strengths", strengths)
            with col3:
                # Count weaknesses (scores < 70)
                weaknesses = sum(1 for component in results.values() if component["score"] < 70)
                st.metric("Areas for Improvement", weaknesses)
            
            # Create a dataframe for visualization
            df = pd.DataFrame({
                "Component": [component.replace('_', ' ').title() for component in results.keys()],
                "Score": [component["score"] for component in results.values()]
            })
            
            # Sort by score
            df = df.sort_values("Score", ascending=False)
            
            # Create two columns for charts
            col1, col2 = st.columns(2)
            
            with col1:
                st.subheader("Component Scores")
                # Create a horizontal bar chart
                fig, ax = plt.subplots(figsize=(10, 8))
                bars = ax.barh(df["Component"], df["Score"], color=df["Score"].apply(
                    lambda x: 'green' if x >= 80 else ('orange' if x >= 70 else 'red')
                ))
                
                # Add score labels
                for bar in bars:
                    ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height()/2, 
                            f"{bar.get_width():.0f}", va='center')
                
                ax.set_xlim(0, 100)
                ax.set_xlabel("Score")
                ax.set_title("Business Model Canvas Component Scores")
                st.pyplot(fig)
            
            with col2:
                st.subheader("Radar Chart")
                # Prepare data for radar chart
                radar_df = pd.DataFrame({
                    "Component": [component.replace('_', ' ').title() for component in results.keys() if component != "overview"],
                    "Score": [component["score"] for key, component in results.items() if key != "overview"]
                })
                
                # Create radar chart using matplotlib
                categories = radar_df["Component"].tolist()
                values = radar_df["Score"].tolist()
                
                # Close the circle by appending the first value to the end
                categories = categories + [categories[0]]
                values = values + [values[0]]
                
                # Calculate angles for the radar chart
                angles = [n / float(len(categories)-1) * 2 * 3.14159 for n in range(len(categories))]
                
                # Create figure
                fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
                
                # Draw one axis per variable and add labels
                plt.xticks(angles[:-1], categories[:-1], size=8)
                
                # Draw the chart
                ax.plot(angles, values, linewidth=2, linestyle='solid')
                ax.fill(angles, values, alpha=0.25)
                
                # Set y-axis limits
                ax.set_ylim(0, 100)
                
                # Add a title
                plt.title("Business Model Canvas Radar", size=14)
                
                # Show the plot
                st.pyplot(fig)
            
            # Key insights and recommendations
            st.subheader("Key Insights")
            
            # Display strengths
            st.markdown("#### Strengths")
            for component, details in sorted(results.items(), key=lambda x: x[1]["score"], reverse=True):
                if details["score"] >= 80:
                    st.markdown(f"**{component.replace('_', ' ').title()} ({details['score']}/100)**: {details['feedback'][:100]}...")
            
            # Display weaknesses
            st.markdown("#### Areas for Improvement")
            for component, details in sorted(results.items(), key=lambda x: x[1]["score"]):
                if details["score"] < 70:
                    st.markdown(f"**{component.replace('_', ' ').title()} ({details['score']}/100)**: {details['feedback'][:100]}...")
        
        with tab2:
            # Create a selectbox for component selection
            component_names = {comp: comp.replace('_', ' ').title() for comp in results.keys()}
            selected_component = st.selectbox(
                "Select a component to view detailed feedback:",
                options=list(results.keys()),
                format_func=lambda x: component_names[x]
            )
            
            # Display the selected component's details
            if selected_component:
                component_data = results[selected_component]
                
                # Create columns for score and feedback
                col1, col2 = st.columns([1, 3])
                
                with col1:
                    # Display score with color
                    score = component_data["score"]
                    score_color = "green" if score >= 80 else ("orange" if score >= 70 else "red")
                    st.markdown(f"### Score: <span style='color:{score_color}'>{score}/100</span>", unsafe_allow_html=True)
                    
                    # Add a visual gauge
                    fig, ax = plt.subplots(figsize=(3, 3))
                    ax.add_patch(plt.Circle((0.5, 0.5), 0.4, color='lightgrey'))
                    ax.add_patch(plt.Circle((0.5, 0.5), 0.4, color=score_color, 
                                          theta1=0, theta2=3.6 * score, 
                                          fill=False, linewidth=10))
                    ax.text(0.5, 0.5, f"{score}", horizontalalignment='center', 
                           verticalalignment='center', fontsize=20)
                    ax.set_xlim(0, 1)
                    ax.set_ylim(0, 1)
                    ax.set_aspect('equal')
                    ax.axis('off')
                    st.pyplot(fig)
                    
                with col2:
                    st.markdown("### Feedback")
                    st.write(component_data["feedback"])
                    
                    # Add improvement priority
                    priority = "Low" if score >= 80 else ("Medium" if score >= 70 else "High")
                    priority_color = "green" if priority == "Low" else ("orange" if priority == "Medium" else "red")
                    st.markdown(f"**Improvement Priority**: <span style='color:{priority_color}'>{priority}</span>", unsafe_allow_html=True)
    
    # Add sidebar with instructions and info
    with st.sidebar:
        st.header("About")
        st.info("""
        This tool uses the Anthropic Claude API to analyze your pitch deck against the Business Model Canvas framework.
        
        Upload your pitch deck (PDF, PPTX, or DOCX) and get detailed feedback on each component of your business model.
        
        The tool evaluates:
        - Key Partners
        - Key Activities
        - Key Resources
        - Value Proposition
        - Customer Relationships
        - Channels
        - Customer Segments
        - Cost Structure
        - Revenue Streams
        """)
        
        st.header("How to Use")
        st.markdown("""
        1. Upload your pitch deck using the file uploader
        2. Click "Analyze Pitch Deck" to start the analysis
        3. View the overview dashboard for a quick assessment
        4. Explore detailed feedback for each component
        """)
        
        # Add export functionality
        if hasattr(st.session_state, 'results') and st.session_state.results:
            st.header("Export Results")
            
            # Create a JSON string of the results
            results_json = json.dumps(st.session_state.results, indent=2)
            
            # Add a download button
            st.download_button(
                label="Download Results (JSON)",
                data=results_json,
                file_name="pitch_deck_assessment.json",
                mime="application/json"
            )
            
            # Also offer CSV export for the scores
            if st.session_state.results:
                csv_data = pd.DataFrame({
                    "Component": [component.replace('_', ' ').title() for component in st.session_state.results.keys()],
                    "Score": [component["score"] for component in st.session_state.results.values()],
                    "Feedback": [component["feedback"] for component in st.session_state.results.values()]
                })
                
                st.download_button(
                    label="Download Results (CSV)",
                    data=csv_data.to_csv(index=False),
                    file_name="pitch_deck_assessment.csv",
                    mime="text/csv"
                )

if __name__ == "__main__":
    main()